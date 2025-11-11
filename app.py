from flask import Flask, request, jsonify
from flask_cors import CORS
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import io
from transformers import pipeline
from pptx import Presentation
import docx
import os
from werkzeug.utils import secure_filename

app = Flask(__name__)
CORS(app)

# ✅ Health check route (Render will use this)
@app.route("/health")
def health():
    return jsonify({"status": "ok"}), 200


# ✅ Load Hugging Face summarizer
try:
    summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
except Exception as e:
    print("⚠️ Failed to load summarizer model:", e)
    summarizer = None


# ---------- Utility Functions ----------

def chunk_text(text, max_len=1500):
    """Split large text into smaller chunks for summarization."""
    words = text.split()
    for i in range(0, len(words), max_len):
        yield ' '.join(words[i:i + max_len])


def summarize_text(text):
    """Summarize long text using chunked summarization."""
    if not summarizer:
        return "Summarizer model not available."

    summaries = []
    for chunk in chunk_text(text):
        if len(chunk.split()) < 20:
            continue
        try:
            result = summarizer(chunk, max_length=150, min_length=40, do_sample=False)
            summaries.append(result[0]['summary_text'])
        except Exception as e:
            print("⚠️ Error summarizing chunk:", e)
    return " ".join(summaries) if summaries else "No significant text found."


def clean_text(text):
    """Remove extra whitespace and empty lines."""
    return ' '.join(text.split())


# ---------- Summarization Endpoint ----------

@app.route('/summarize', methods=['POST'])
def summarize():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    filename = secure_filename(file.filename)
    ext = filename.split('.')[-1].lower()

    results = {}

    try:
        # ---- PDF Handling ----
        if ext == "pdf":
            doc = fitz.open(stream=file.read(), filetype="pdf")
            pages = []
            print(f"📘 Processing {len(doc)} PDF pages...")
            for i, page in enumerate(doc):
                text = clean_text(page.get_text("text"))
                if len(text.strip()) > 5:
                    pages.append((i + 1, text))
            results = {f"Page {num}": summarize_text(text) for num, text in sorted(pages, key=lambda x: x[0])}

        # ---- Image Handling ----
        elif ext in ["jpg", "jpeg", "png"]:
            image = Image.open(io.BytesIO(file.read()))
            text = clean_text(pytesseract.image_to_string(image))
            results["Summary"] = summarize_text(text) if text.strip() else "No readable text found in image."

        # ---- PowerPoint Handling ----
        elif ext == "pptx":
            prs = Presentation(io.BytesIO(file.read()))
            slides = []
            for i, slide in enumerate(prs.slides):
                text = " ".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
                text = clean_text(text)
                if len(text.strip()) > 5:
                    slides.append((i + 1, text))
            results = {f"Slide {num}": summarize_text(text) for num, text in sorted(slides, key=lambda x: x[0])}

        # ---- Word Document Handling ----
        elif ext == "docx":
            doc = docx.Document(io.BytesIO(file.read()))
            sections = [(i + 1, clean_text(p.text)) for i, p in enumerate(doc.paragraphs) if len(p.text.strip()) > 5]
            results = {f"Section {num}": summarize_text(text) for num, text in sorted(sections, key=lambda x: x[0])}

        else:
            return jsonify({"error": "Unsupported file type"}), 400

        return jsonify({"summaries": results})

    except Exception as e:
        print("❌ Error processing file:", e)
        return jsonify({"error": str(e)}), 500


# ---------- Home Route ----------

@app.route('/')
def home():
    return jsonify({
        "message": "🚀 DocVibe AI Summarizer API is live on Render!",
        "endpoints": {
            "/health": "Health check",
            "/summarize": "POST endpoint for file summarization"
        }
    })


# ---------- Run Flask ----------

if __name__ == '__main__':
    port = int(os.environ.get("PORT", 5000))
    print(f"🚀 Starting Flask server on port {port}")
    app.run(host='0.0.0.0', port=port)
